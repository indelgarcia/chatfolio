"""
ChatFolio Final Presentation Generator
Run from project root: python evaluation/generate_final_slides.py
Outputs: evaluation/chatfolio_final_presentation.pptx
"""
import io
import os

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ────────────────────────────────────────────────────────────────────
C_BG    = RGBColor(26,  26,  46)
C_CARD  = RGBColor(30,  41,  59)
C_TEAL  = RGBColor(78,  205, 196)
C_RED   = RGBColor(255, 107, 107)
C_WHITE = RGBColor(255, 255, 255)
C_LGRAY = RGBColor(200, 200, 200)
C_DGRAY = RGBColor(120, 120, 140)
C_GREEN = RGBColor(34,  197, 94)
C_YELL  = RGBColor(255, 213, 79)
C_DARK  = RGBColor(13,  17,  23)
C_ORANGE = RGBColor(249, 115, 22)

THEME_C = [
    RGBColor(239, 68,  68),
    RGBColor(249, 115, 22),
    RGBColor(234, 179, 8),
    RGBColor(34,  197, 94),
    RGBColor(59,  130, 246),
    RGBColor(139, 92,  246),
]

W = Inches(13.33)
H = Inches(7.5)

CATS = ['Goal', 'Timeline', 'Budget', 'Risk', 'Additional', 'Post-Gen', 'Off-Topic']
N    = [7, 7, 7, 6, 7, 7, 8]
V1_H = [71.4, 71.4, 71.4, 66.7, 28.6, 85.7, 100.0]
V2_H = [100.0, 85.7, 100.0, 100.0, 100.0, 85.7, 100.0]

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


# ── pptx helpers ──────────────────────────────────────────────────────────────
def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return slide


def txt(slide, text, l, t, w, h, sz=16, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(sz)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = 'Calibri'
    return box


def bullets(slide, items, l, t, w, h, sz=14, color=C_LGRAY, prefix='  •  '):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text           = prefix + item
        run.font.size      = Pt(sz)
        run.font.color.rgb = color
        run.font.name      = 'Calibri'
    return box


def rect(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp


def img(slide, src, l, t, w, h=None):
    if isinstance(src, str):
        return slide.shapes.add_picture(src, l, t, width=w, height=h) if h \
               else slide.shapes.add_picture(src, l, t, width=w)
    src.seek(0)
    return slide.shapes.add_picture(src, l, t, width=w, height=h) if h \
           else slide.shapes.add_picture(src, l, t, width=w)


def accent(slide):
    rect(slide, 0, 0, W, Inches(0.07), C_TEAL)


def header(slide, title, tag=''):
    accent(slide)
    txt(slide, title, Inches(0.5), Inches(0.1), Inches(10.5), Inches(0.72),
        sz=26, bold=True)
    if tag:
        txt(slide, tag, Inches(10.0), Inches(0.1), Inches(3.0), Inches(0.5),
            sz=10, color=C_TEAL, align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def mkchart(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=150,
                facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf


# ── Chart ─────────────────────────────────────────────────────────────────────
def comparison_chart():
    fig, ax = plt.subplots(figsize=(8.0, 3.8), facecolor='#0D1117')
    ax.set_facecolor('#0D1117')
    y   = np.arange(len(CATS))
    bh  = 0.35
    ax.barh(y - bh/2, V1_H, bh, color='#FF6B6B', label='V1 Baseline', zorder=3)
    b2 = ax.barh(y + bh/2, V2_H, bh, color='#4ECDC4', label='V2 Refined',  zorder=3)
    ax.set_yticks(y)
    ax.set_yticklabels(CATS, color='white', fontsize=11)
    ax.set_xlabel('Human Pass Rate (%)', color='#AAAAAA', fontsize=10)
    ax.set_xlim(0, 125)
    ax.tick_params(colors='#AAAAAA')
    for sp in ['top', 'right']:   ax.spines[sp].set_visible(False)
    for sp in ['bottom', 'left']: ax.spines[sp].set_color('#444')
    ax.grid(axis='x', color='#2A2A3E', linewidth=0.6, zorder=0)
    for b, v in zip(b2, V2_H):
        ax.text(v + 1.5, b.get_y() + b.get_height() / 2, f'{v:.0f}%',
                va='center', color='#4ECDC4', fontsize=9)
    ax.legend(loc='lower right', facecolor='#1E293B', labelcolor='white', fontsize=10)
    ax.set_title('Before vs. After: Human Pass Rate by Category', color='white',
                 fontsize=12, pad=8)
    fig.tight_layout()
    return mkchart(fig)


# ── Slides ────────────────────────────────────────────────────────────────────
def s01_title(prs):
    s = new_slide(prs)
    rect(s, 0, 0, W, Inches(0.08), C_TEAL)
    txt(s, 'ChatFolio', Inches(1), Inches(1.2), Inches(11.3), Inches(1.4),
        sz=64, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 'Final Presentation', Inches(1), Inches(2.6), Inches(11.3), Inches(0.8),
        sz=32, color=C_TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'AI-Powered Conversational Investment Advisor',
        Inches(1), Inches(3.4), Inches(11.3), Inches(0.6),
        sz=18, color=C_LGRAY, align=PP_ALIGN.CENTER)
    rect(s, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.03), C_DGRAY)
    txt(s, 'Chris  ·  Indel  ·  Nate  ·  Matt     |     April 2026',
        Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
        sz=14, color=C_DGRAY, align=PP_ALIGN.CENTER)
    notes(s, (
        "Welcome everyone. Today we're presenting ChatFolio, our HAI term project. "
        "We set out to solve a specific problem: why do so many people who want to "
        "invest never actually start? We'll walk through our design journey, "
        "what we built, how we evaluated it, and what we learned."
    ))


def s02_problem(prs):
    s = new_slide(prs)
    header(s, 'The Problem', 'PROJECT OVERVIEW')
    # Stat cards left
    for i, (stat, label) in enumerate([
        ('57%', "of US adults don't personally own stocks"),
        ('38%', "of non-owners say they don't understand the market"),
    ]):
        y = Inches(1.1) + i * Inches(2.1)
        rect(s, Inches(0.4), y, Inches(4.8), Inches(1.8), C_CARD)
        rect(s, Inches(0.4), y, Inches(0.12), Inches(1.8), C_RED)
        txt(s, stat, Inches(0.7), y + Inches(0.1), Inches(4.3), Inches(0.85),
            sz=48, bold=True, color=C_RED)
        txt(s, label, Inches(0.7), y + Inches(0.95), Inches(4.3), Inches(0.75),
            sz=13, color=C_LGRAY)
    txt(s, 'Source: Philadelphia Fed, Consumer Finance Brief',
        Inches(0.4), Inches(5.45), Inches(5.0), Inches(0.35),
        sz=9, color=C_DGRAY, italic=True)
    # Pain points right
    rect(s, Inches(5.8), Inches(1.1), Inches(7.1), Inches(5.9), C_CARD)
    txt(s, 'WHY BEGINNERS DON\'T INVEST', Inches(6.05), Inches(1.2), Inches(6.6), Inches(0.45),
        sz=13, bold=True, color=C_RED)
    pain = [
        'Information overload — too many products, tickers, and strategies',
        'Fear of making a costly mistake with real money',
        'Brokerage forms require financial knowledge to complete',
        'Existing tools tell you WHAT to buy — not WHY it fits you',
        'No clear starting point for someone with $100–$200/month',
    ]
    bullets(s, pain, Inches(6.05), Inches(1.75), Inches(6.6), Inches(4.8),
            sz=14, prefix='  →  ')
    notes(s, (
        "The problem is real and measurable. More than half of American adults have "
        "no stock market exposure, and a big driver is a lack of understanding — not "
        "a lack of desire. Existing tools like Robinhood and Fidelity are built for "
        "people who already know what they want. ChatFolio is built for people who don't."
    ))


def s03_solution(prs):
    s = new_slide(prs)
    header(s, 'Our Solution + Target User', 'PROJECT OVERVIEW')
    # Solution left
    rect(s, Inches(0.4), Inches(1.0), Inches(6.5), Inches(5.9), C_CARD)
    txt(s, 'CHATFOLIO', Inches(0.65), Inches(1.1), Inches(6.0), Inches(0.45),
        sz=13, bold=True, color=C_TEAL)
    txt(s, 'A conversational investment advisor that replaces brokerage forms with natural dialogue.',
        Inches(0.65), Inches(1.6), Inches(6.0), Inches(0.75),
        sz=14, color=C_LGRAY)
    flow = [
        'User describes goals in plain language',
        'GPT-4o-mini extracts profile through conversation',
        'Live sidebar shows what the system has captured',
        'User clicks "Generate My Portfolio" (human-in-the-loop)',
        'Rule-based ETF allocation + AI-written rationale',
        'Growth projections + step-by-step action plan',
        'User can chat to adjust — portfolio regenerates',
    ]
    bullets(s, flow, Inches(0.65), Inches(2.45), Inches(6.0), Inches(4.2),
            sz=13, color=C_WHITE, prefix='  ✦  ')
    # Marcus card right
    rect(s, Inches(7.2), Inches(1.0), Inches(5.7), Inches(5.9), C_CARD)
    rect(s, Inches(7.2), Inches(1.0), Inches(5.7), Inches(0.08), C_TEAL)
    txt(s, 'MEET MARCUS — OUR TARGET USER', Inches(7.45), Inches(1.12), Inches(5.2), Inches(0.4),
        sz=11, bold=True, color=C_TEAL)
    attrs = [
        ('Age', '26 — recent grad, first full-time job'),
        ('Budget', '~$150/month to invest'),
        ('Knowledge', 'Intimidated by financial jargon'),
        ('Comfort', 'Comfortable with chat interfaces'),
        ('Goal', 'Build wealth without learning finance theory'),
        ('Problem', 'Doesn\'t know where or how to start'),
    ]
    for i, (k, v) in enumerate(attrs):
        y = Inches(1.65) + i * Inches(0.75)
        txt(s, k, Inches(7.45), y, Inches(1.4), Inches(0.5),
            sz=12, bold=True, color=C_TEAL)
        txt(s, v, Inches(8.9),  y, Inches(3.8), Inches(0.5),
            sz=12, color=C_LGRAY)
    notes(s, (
        "ChatFolio's core idea: replace the intimidating brokerage form with a conversation. "
        "Marcus represents our target user — someone who wants to invest but keeps hitting "
        "a wall of jargon and complexity. Every design decision was filtered through the question: "
        "would Marcus understand this? Would Marcus trust this?"
    ))


def s04_how_it_works(prs):
    s = new_slide(prs)
    header(s, 'How It Works', 'PROJECT OVERVIEW')
    # Flow nodes
    nodes = [
        ('User\nChat', C_TEAL),
        ('GPT-4o-mini\n(JSON mode)', RGBColor(99, 102, 241)),
        ('Live Profile\nSidebar', RGBColor(59, 130, 246)),
        ('Generate\nButton', C_ORANGE),
        ('Portfolio +\nRationale +\nProjections', C_GREEN),
    ]
    nw, nh = Inches(2.0), Inches(1.1)
    gap    = Inches(0.32)
    total  = len(nodes) * nw + (len(nodes) - 1) * gap
    sx     = (W - total) / 2
    sy     = Inches(1.5)
    for i, (label, color) in enumerate(nodes):
        x = sx + i * (nw + gap)
        rect(s, x, sy, nw, nh, color)
        txt(s, label, x, sy, nw, nh, sz=13, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            ax = x + nw + Inches(0.04)
            ay = sy + nh / 2 - Inches(0.05)
            txt(s, '→', ax, ay, Inches(0.22), Inches(0.35),
                sz=16, bold=True, color=C_LGRAY, align=PP_ALIGN.CENTER)
    # HAI criteria row
    txt(s, 'HAI Evaluation Criteria:', Inches(0.4), Inches(3.0), Inches(4.0), Inches(0.4),
        sz=13, bold=True, color=C_LGRAY)
    criteria = [
        ('Onboarding', RGBColor(239, 68, 68),   'Explainer panel, disclaimer,\npost-gen guidance banner'),
        ('User Input',  RGBColor(249, 115, 22),  'Conversational Q&A, one question\nat a time, open-ended extras'),
        ('System Output', RGBColor(34, 197, 94), 'Live sidebar, portfolio rationale,\nprojections, action plan'),
        ('User Feedback', RGBColor(59, 130, 246),'Confirmation loops, edit loop,\nspinner + toast on update'),
    ]
    cw = Inches(3.1)
    for i, (name, color, desc) in enumerate(criteria):
        x = Inches(0.4) + i * (cw + Inches(0.12))
        rect(s, x, Inches(3.55), cw, Inches(3.5), C_CARD)
        rect(s, x, Inches(3.55), cw, Inches(0.08), color)
        txt(s, name, x + Inches(0.15), Inches(3.67), cw - Inches(0.3), Inches(0.45),
            sz=13, bold=True, color=color)
        txt(s, desc, x + Inches(0.15), Inches(4.18), cw - Inches(0.3), Inches(2.7),
            sz=11, color=C_LGRAY)
    notes(s, (
        "Here's the full pipeline. Every user message flows through GPT-4o-mini in JSON mode — "
        "the model returns both a conversational reply and structured profile updates in one call. "
        "The generate button is deliberate friction: the user confirms their profile before we "
        "produce a recommendation. Below, you can see how each part of the system maps to the "
        "four HAI evaluation criteria — this wasn't an afterthought, it was the design constraint."
    ))


def s05_v01_prototype(prs):
    s = new_slide(prs)
    header(s, 'v0.1 — Initial Prototype', 'PROTOTYPE EVOLUTION  |  March 2026')
    txt(s, 'Lovable no-code mockup — static UI, no real AI',
        Inches(0.5), Inches(0.82), Inches(9.0), Inches(0.35),
        sz=13, color=C_DGRAY, italic=True)
    # Image frames
    img_h = Inches(4.5)
    img_w = Inches(5.8)
    # Frame rects (slightly larger than images)
    rect(s, Inches(0.35), Inches(1.25), img_w + Inches(0.1), img_h + Inches(0.1),
         RGBColor(40, 50, 70))
    rect(s, Inches(6.65), Inches(1.25), img_w + Inches(0.1), img_h + Inches(0.1),
         RGBColor(40, 50, 70))
    # Embed images
    lp = os.path.join(ROOT, 'context', 'loveable_1st_prototype', 'loveable_landing_page.png')
    gp = os.path.join(ROOT, 'context', 'loveable_1st_prototype', 'loveable_goal.png')
    img(s, lp, Inches(0.4),  Inches(1.3),  img_w, img_h)
    img(s, gp, Inches(6.7),  Inches(1.3),  img_w, img_h)
    txt(s, 'Landing Page', Inches(0.4),  Inches(5.82), img_w, Inches(0.3),
        sz=10, color=C_DGRAY, align=PP_ALIGN.CENTER)
    txt(s, 'Chat + Live Profile Sidebar', Inches(6.7), Inches(5.82), img_w, Inches(0.3),
        sz=10, color=C_DGRAY, align=PP_ALIGN.CENTER)
    # Callout pills
    callouts = [
        ('✓  Established the dual-panel layout concept', C_GREEN),
        ('✓  Proved conversational framing over forms', C_GREEN),
        ('✗  Static — no real AI or profile extraction', C_RED),
    ]
    for i, (label, color) in enumerate(callouts):
        x = Inches(0.4) + i * Inches(4.3)
        rect(s, x, Inches(6.2), Inches(4.1), Inches(0.62), C_CARD)
        txt(s, label, x + Inches(0.12), Inches(6.25), Inches(3.9), Inches(0.5),
            sz=12, color=color)
    notes(s, (
        "Before writing a single line of Python, we built a static mockup in Lovable. "
        "This let the four of us align on the visual concept and interaction model in hours, "
        "not days. The dual-panel layout — chat on the left, live profile on the right — "
        "came directly from this mockup. The Lovable version even had Accept/Adjust/Explain buttons "
        "on the portfolio, which set the aspiration for what the real system should eventually do."
    ))


def s06_v02_prototype(prs):
    s = new_slide(prs)
    header(s, 'v0.2 / v0.3 — First Functional Prototype', 'PROTOTYPE EVOLUTION  |  March 2026')
    txt(s, 'Streamlit + GPT-4o-mini — real AI, real profile extraction',
        Inches(0.5), Inches(0.82), Inches(9.0), Inches(0.35),
        sz=13, color=C_DGRAY, italic=True)

    # Feature list left
    rect(s, Inches(0.4), Inches(1.25), Inches(5.0), Inches(5.85), C_CARD)
    txt(s, 'v0.2  —  CORE ENGINE', Inches(0.65), Inches(1.35), Inches(4.6), Inches(0.42),
        sz=12, bold=True, color=C_TEAL)
    bullets(s, [
        'GPT-4o-mini live with JSON response mode',
        'Conversational profile extraction',
        'Confirmation loops on ambiguous answers',
        'Live sidebar updates in real time',
        '"Generate My Portfolio" button (human-in-the-loop)',
        'Rule-based ETF allocation (VTI/VXUS/BND/BNDX)',
    ], Inches(0.65), Inches(1.82), Inches(4.6), Inches(2.5), sz=12)
    txt(s, 'v0.3  —  FULL FEATURE SET', Inches(0.65), Inches(4.42), Inches(4.6), Inches(0.42),
        sz=12, bold=True, color=C_TEAL)
    bullets(s, [
        '"What is ChatFolio?" onboarding explainer panel',
        'AI-generated portfolio rationale (plain language)',
        'Growth projections — 3 scenarios, line chart',
        'Step-by-step action plan (brokerage, account, ETFs)',
        'Post-generation edit loop — chat stays active',
        'Static disclaimer + guardrails in system prompt',
    ], Inches(0.65), Inches(4.9), Inches(4.6), Inches(2.0), sz=12)

    # Wireframe right
    wf_x = Inches(5.7)
    wf_w = Inches(7.2)
    wf_h = Inches(5.85)
    rect(s, wf_x, Inches(1.25), wf_w, wf_h, C_DARK)

    # Sidebar strip
    sb_w = Inches(1.7)
    rect(s, wf_x, Inches(1.25), sb_w, wf_h, C_CARD)
    txt(s, 'Your Profile', wf_x + Inches(0.08), Inches(1.32), sb_w - Inches(0.1), Inches(0.3),
        sz=9, bold=True, color=C_TEAL)
    for i, (field, val) in enumerate([
        ('Goal', 'Retirement'), ('Timeline', '30 years'),
        ('Budget', '$150/mo'), ('Risk', 'Moderate'), ('Extra', '401k noted'),
    ]):
        fy = Inches(1.7) + i * Inches(0.68)
        rect(s, wf_x + Inches(0.08), fy, sb_w - Inches(0.16), Inches(0.58),
             RGBColor(20, 30, 50))
        txt(s, field, wf_x + Inches(0.14), fy + Inches(0.02), sb_w - Inches(0.22), Inches(0.25),
            sz=7, color=C_DGRAY)
        txt(s, val,   wf_x + Inches(0.14), fy + Inches(0.26), sb_w - Inches(0.22), Inches(0.25),
            sz=9, bold=True, color=C_WHITE)
    txt(s, 'Progress  4/4', wf_x + Inches(0.08), Inches(5.25), sb_w - Inches(0.1), Inches(0.25),
        sz=8, color=C_GREEN)
    rect(s, wf_x + Inches(0.08), Inches(5.5), sb_w - Inches(0.16), Inches(0.1), C_GREEN)

    # Chat area
    cx = wf_x + sb_w + Inches(0.08)
    cw = wf_w - sb_w - Inches(0.16)
    bubbles = [
        ('Hi! I\'m ChatFolio. What are you investing for?', True),
        ('Saving for retirement, about 30 years out', False),
        ('Great! How much can you invest each month?', True),
        ('Around $150', False),
    ]
    for i, (msg, is_ai) in enumerate(bubbles):
        bw = cw * 0.72
        bx = cx if is_ai else cx + cw - bw - Inches(0.05)
        by = Inches(1.35) + i * Inches(0.82)
        bc = RGBColor(30, 41, 59) if is_ai else RGBColor(20, 60, 80)
        tc = C_TEAL if is_ai else C_WHITE
        rect(s, bx, by, bw, Inches(0.65), bc)
        txt(s, msg, bx + Inches(0.08), by + Inches(0.06), bw - Inches(0.16), Inches(0.55),
            sz=9, color=tc)

    # Generate button
    btn_y = Inches(5.5)
    rect(s, cx, btn_y, cw - Inches(0.05), Inches(0.42), C_TEAL)
    txt(s, 'Generate My Portfolio  →', cx, btn_y, cw - Inches(0.05), Inches(0.42),
        sz=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # Input bar
    rect(s, cx, Inches(6.0), cw - Inches(0.05), Inches(0.35), RGBColor(20, 28, 42))
    txt(s, 'Type your message...', cx + Inches(0.08), Inches(6.02),
        cw - Inches(0.25), Inches(0.3), sz=9, color=C_DGRAY, italic=True)

    txt(s, 'Wireframe — Streamlit dark theme', wf_x, Inches(7.15), wf_w, Inches(0.28),
        sz=8, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)

    notes(s, (
        "In v0.2 we replaced the static mockup with a real Streamlit app backed by GPT-4o-mini. "
        "Users could now have an actual conversation and see their profile fill in live. "
        "v0.3 added everything that makes the output meaningful: the rationale that explains WHY "
        "the portfolio was built this way, the growth projections that make the numbers concrete, "
        "and the action plan that tells the user what to actually do next. "
        "The edit loop was also added here — users could keep chatting after generation."
    ))


def s07_v05_prototype(prs):
    s = new_slide(prs)
    header(s, 'v0.4 / v0.5 — Post-Evaluation Polish', 'PROTOTYPE EVOLUTION  |  April 2026')
    txt(s, 'Prompt refinements + UX improvements driven by user testing',
        Inches(0.5), Inches(0.82), Inches(9.0), Inches(0.35),
        sz=13, color=C_DGRAY, italic=True)

    # Changes left
    rect(s, Inches(0.4), Inches(1.25), Inches(5.4), Inches(5.85), C_CARD)
    txt(s, 'PROMPT CHANGES (v0.4)', Inches(0.65), Inches(1.35), Inches(5.0), Inches(0.4),
        sz=12, bold=True, color=C_TEAL)
    p_changes = [
        ('1', 'One-field confirmation loop — paraphrase + confirm before saving'),
        ('2', 'Edge case handling: variable budget, $0 budget, short timeline'),
        ('3', 'Post-gen behavioral split: adjust profile vs. answer questions'),
        ('4', 'Proactive contextual question after 4 required fields'),
        ('5', 'Plain-language risk descriptions required — no numeric confidence'),
    ]
    for i, (num, desc) in enumerate(p_changes):
        y = Inches(1.85) + i * Inches(0.7)
        rect(s, Inches(0.65), y, Inches(0.35), Inches(0.35), C_TEAL)
        txt(s, num, Inches(0.65), y, Inches(0.35), Inches(0.35),
            sz=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(1.1), y, Inches(4.5), Inches(0.5),
            sz=11, color=C_LGRAY)

    txt(s, 'UX CHANGES (v0.5)', Inches(0.65), Inches(5.4), Inches(5.0), Inches(0.4),
        sz=12, bold=True, color=C_TEAL)
    bullets(s, [
        'Split view: chat and portfolio scroll independently',
        'Spinner + success toast on portfolio regeneration',
        'Prominent follow-up guidance banner after generation',
    ], Inches(0.65), Inches(5.85), Inches(5.0), Inches(1.1), sz=11)

    # Split-view wireframe right
    wf_x = Inches(6.05)
    wf_w = Inches(6.85)
    wf_h = Inches(5.85)
    rect(s, wf_x, Inches(1.25), wf_w, wf_h, C_DARK)

    # Left chat column
    lc_w = wf_w * 0.4
    rect(s, wf_x, Inches(1.25), lc_w, wf_h, C_CARD)
    txt(s, 'Chat', wf_x + Inches(0.1), Inches(1.32), lc_w - Inches(0.2), Inches(0.3),
        sz=9, bold=True, color=C_TEAL)
    for i, (msg, is_ai) in enumerate([
        ('Portfolio generated!', True),
        ('Make it more aggressive', False),
        ('Updated to Aggressive. Regenerating...', True),
    ]):
        by = Inches(1.72) + i * Inches(0.72)
        bc = RGBColor(20, 35, 55) if is_ai else RGBColor(20, 55, 75)
        tc = C_TEAL if is_ai else C_WHITE
        rect(s, wf_x + Inches(0.08), by, lc_w - Inches(0.16), Inches(0.58), bc)
        txt(s, msg, wf_x + Inches(0.14), by + Inches(0.06), lc_w - Inches(0.28), Inches(0.46),
            sz=8, color=tc)
    # Banner
    rect(s, wf_x, Inches(5.88), lc_w, Inches(0.5), RGBColor(10, 40, 60))
    txt(s, '💡 Try: "make it more aggressive" or "what is VTI?"',
        wf_x + Inches(0.08), Inches(5.9), lc_w - Inches(0.1), Inches(0.45),
        sz=7, color=C_TEAL)

    # Right portfolio column
    rc_x = wf_x + lc_w + Inches(0.05)
    rc_w = wf_w - lc_w - Inches(0.05)
    rect(s, rc_x, Inches(1.25), rc_w, wf_h, RGBColor(15, 22, 35))

    sections = [
        ('Allocation', [('VTI', 0.64, C_TEAL), ('VXUS', 0.21, C_GREEN),
                        ('BND', 0.1, C_YELL), ('BNDX', 0.05, C_ORANGE)]),
    ]
    txt(s, 'ALLOCATION', rc_x + Inches(0.1), Inches(1.35), rc_w, Inches(0.3),
        sz=8, bold=True, color=C_DGRAY)
    for i, (ticker, pct, color) in enumerate([
        ('VTI', 0.64, C_TEAL), ('VXUS', 0.21, C_GREEN), ('BND', 0.10, C_YELL), ('BNDX', 0.05, C_ORANGE)
    ]):
        by = Inches(1.72) + i * Inches(0.52)
        bar_w = (rc_w - Inches(0.2)) * pct
        txt(s, ticker, rc_x + Inches(0.1), by, Inches(0.5), Inches(0.25),
            sz=8, bold=True, color=color)
        rect(s, rc_x + Inches(0.1), by + Inches(0.27), bar_w, Inches(0.18), color)
        txt(s, f'{int(pct*100)}%', rc_x + Inches(0.1) + bar_w + Inches(0.05), by + Inches(0.27),
            Inches(0.4), Inches(0.2), sz=7, color=C_LGRAY)

    for label, y in [('RATIONALE', 4.0), ('PROJECTIONS', 4.7), ('ACTION PLAN', 5.4)]:
        rect(s, rc_x + Inches(0.1), Inches(y - 0.05), rc_w - Inches(0.2), Inches(0.55),
             RGBColor(20, 30, 48))
        txt(s, label, rc_x + Inches(0.18), Inches(y), rc_w - Inches(0.36), Inches(0.3),
            sz=8, bold=True, color=C_DGRAY)

    txt(s, 'Wireframe — post-generation split view', wf_x, Inches(7.15), wf_w, Inches(0.28),
        sz=8, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)

    notes(s, (
        "After running our think-aloud study, we had a concrete list of things to fix. "
        "v0.4 was a prompt overhaul: we tightened the confirmation loop, handled edge cases "
        "we'd seen break in testing, and required the AI to explain risk in plain English every time. "
        "v0.5 was the UX pass: the post-generation layout now uses two independent scroll containers "
        "so users can review the portfolio while keeping the conversation in view. "
        "The spinner and success toast give clear feedback that a follow-up message triggered a change."
    ))


def s08_study(prs):
    s = new_slide(prs)
    header(s, 'Evaluation: Think-Aloud Study', 'EVALUATION & REFINEMENT  |  n = 7')
    # Method box
    rect(s, Inches(0.4), Inches(1.0), Inches(3.5), Inches(5.9), C_CARD)
    txt(s, 'METHODOLOGY', Inches(0.65), Inches(1.1), Inches(3.0), Inches(0.4),
        sz=11, bold=True, color=C_TEAL)
    bullets(s, [
        'Task: complete full ChatFolio flow',
        '   greeting → 4 questions → generate',
        'Method: think-aloud, moderated',
        'Observer noted friction points and',
        '   surprises without guiding users',
        'n = 7 college students',
        '   Mixed investing experience',
        '   CS, Finance, and Economics majors',
        'Synthesis: affinity diagramming → 6 themes',
    ], Inches(0.65), Inches(1.6), Inches(3.0), Inches(5.0), sz=12, prefix='')

    # 6-card affinity grid
    themes = [
        ('Onboarding &\nInitial Prompt',       '6 / 7 participants', 0),
        ('Risk Tolerance\nConfusion',           '3 / 7 participants', 1),
        ('Chat &\nConversation Flow',           '5 / 7 participants', 2),
        ('Portfolio Output &\nComprehension',   '5 / 7 participants', 3),
        ('UX &\nVisual Polish',                 '4 / 7 participants', 4),
        ('What Worked Well',                    '6 / 7 participants', 5),
    ]
    cw, ch = Inches(3.0), Inches(2.7)
    sx, sy = Inches(4.1), Inches(1.0)
    gx, gy = Inches(0.15), Inches(0.2)
    for i, (name, count, ci) in enumerate(themes):
        col, row = i % 3, i // 3
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)
        tc = THEME_C[ci]
        rect(s, x, y, cw, ch, C_CARD)
        rect(s, x, y, cw, Inches(0.09), tc)
        txt(s, name, x + Inches(0.15), y + Inches(0.18), cw - Inches(0.3), Inches(1.1),
            sz=15, bold=True)
        rect(s, x + Inches(0.15), y + ch - Inches(0.6), Inches(1.5), Inches(0.4), tc)
        txt(s, count, x + Inches(0.15), y + ch - Inches(0.63), Inches(1.5), Inches(0.43),
            sz=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    notes(s, (
        "We ran a think-aloud study with 7 participants who used ChatFolio while narrating their "
        "thoughts aloud. We observed — we didn't guide or help. "
        "After the sessions, we used affinity diagramming to cluster observations into 6 themes. "
        "Onboarding and What Worked Well both hit 6 of 7 participants — "
        "meaning nearly everyone had something to say about the first impression, "
        "and nearly everyone found something they genuinely liked."
    ))


def s09_findings(prs):
    s = new_slide(prs)
    header(s, 'Key Findings → What We Changed', 'EVALUATION & REFINEMENT')

    # Pain points left
    rect(s, Inches(0.4), Inches(1.0), Inches(6.1), Inches(5.9), C_CARD)
    rect(s, Inches(0.4), Inches(1.0), Inches(0.1), Inches(5.9), C_RED)
    txt(s, 'PAIN POINTS OBSERVED', Inches(0.65), Inches(1.08), Inches(5.6), Inches(0.4),
        sz=12, bold=True, color=C_RED)
    pain = [
        'Onboarding prompt implied only "retirement" or "house" as valid goals',
        'Risk tolerance question confused beginners — "35% confident" phrasing broke trust',
        'Users dumped all info in one message, bypassing step-by-step design',
        'ETF tickers (VTI, VXUS, BND) were meaningless — literacy gap',
        'No feedback when a follow-up message triggered portfolio regeneration',
    ]
    for i, p in enumerate(pain):
        y = Inches(1.6) + i * Inches(0.98)
        rect(s, Inches(0.6), y, Inches(5.7), Inches(0.82), RGBColor(40, 20, 20))
        txt(s, f'{i+1}', Inches(0.72), y + Inches(0.1), Inches(0.35), Inches(0.55),
            sz=14, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
        txt(s, p, Inches(1.15), y + Inches(0.1), Inches(5.0), Inches(0.65),
            sz=11, color=C_LGRAY)

    # What we did right
    rect(s, Inches(6.8), Inches(1.0), Inches(6.1), Inches(5.9), C_CARD)
    rect(s, Inches(6.8), Inches(1.0), Inches(0.1), Inches(5.9), C_TEAL)
    txt(s, 'WHAT WE CHANGED', Inches(7.05), Inches(1.08), Inches(5.6), Inches(0.4),
        sz=12, bold=True, color=C_TEAL)
    fixes = [
        'Rewrote opening prompt to accept any goal; removed example anchors',
        'Required plain-language risk descriptions; banned numeric outputs',
        'Enforced one-question-at-a-time in system prompt; added confirmation loop',
        'Added inline ETF explanations to rationale; defined "index fund" on first use',
        'Added spinner + success toast on regeneration; clear visual feedback',
    ]
    for i, f in enumerate(fixes):
        y = Inches(1.6) + i * Inches(0.98)
        rect(s, Inches(7.0), y, Inches(5.7), Inches(0.82), RGBColor(10, 35, 35))
        txt(s, '✓', Inches(7.12), y + Inches(0.1), Inches(0.35), Inches(0.55),
            sz=14, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
        txt(s, f, Inches(7.55), y + Inches(0.1), Inches(5.0), Inches(0.65),
            sz=11, color=C_LGRAY)

    notes(s, (
        "This is the slide that shows why HAI evaluation matters. Each pain point on the left "
        "maps directly to a concrete change on the right. "
        "The financial literacy gap was the biggest surprise — we assumed users would google an ETF ticker. "
        "They didn't. They either guessed or disengaged. "
        "The regeneration feedback issue was simpler but just as impactful: "
        "users thought nothing happened when they sent a follow-up message."
    ))


def s10_results(prs, buf_cmp):
    s = new_slide(prs)
    header(s, 'Evaluation Results', 'EVALUATION & REFINEMENT')
    img(s, buf_cmp, Inches(0.3), Inches(1.0), Inches(8.7))

    # Summary table
    cx = [Inches(9.35), Inches(10.35), Inches(11.35), Inches(12.3)]
    cw = Inches(1.0)
    for j, h in enumerate(['Metric', 'V1', 'V2', 'Δ']):
        txt(s, h, cx[j], Inches(1.08), cw, Inches(0.38),
            sz=12, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    rect(s, Inches(9.35), Inches(1.46), Inches(3.65), Inches(0.03), C_DGRAY)
    rows = [
        ('LLM Avg',      '3.61', '4.80', '+1.19'),
        ('Human Avg',    '3.96', '4.76', '+0.80'),
        ('LLM Pass %',   '67.3%', '91.8%', '+24.5pp'),
        ('Human Pass %', '71.4%', '95.9%', '+24.5pp'),
    ]
    for i, row in enumerate(rows):
        y  = Inches(1.55) + i * Inches(0.65)
        bg = C_CARD if i % 2 == 0 else RGBColor(20, 28, 42)
        rect(s, Inches(9.35), y - Inches(0.05), Inches(3.65), Inches(0.6), bg)
        for j, (cell, cc) in enumerate(zip(row, [C_LGRAY, C_RED, C_TEAL, C_GREEN])):
            txt(s, cell, cx[j], y, cw, Inches(0.5),
                sz=12, bold=(j == 3), color=cc, align=PP_ALIGN.CENTER)

    rect(s, Inches(9.35), Inches(4.3), Inches(3.65), Inches(1.6), RGBColor(10, 30, 30))
    txt(s, '+24.5 pp', Inches(9.35), Inches(4.35), Inches(3.65), Inches(0.7),
        sz=34, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'human pass rate improvement',
        Inches(9.35), Inches(5.05), Inches(3.65), Inches(0.35),
        sz=11, color=C_LGRAY, align=PP_ALIGN.CENTER)
    txt(s, '71.4%  →  95.9%',
        Inches(9.35), Inches(5.42), Inches(3.65), Inches(0.4),
        sz=16, bold=True, align=PP_ALIGN.CENTER)

    notes(s, (
        "The quantitative eval used a 49-sample benchmark dataset across 7 categories. "
        "Each response was rated 1-5 by both an LLM judge and a human reviewer — "
        "pass threshold was 4 or above. "
        "V1 passed 71.4% of human-rated cases. V2 passed 95.9% — a 24.5 percentage point jump. "
        "Crucially, the LLM-judge and human ratings stayed within 0.1 of each other on average, "
        "which gives us confidence the automated metric is actually measuring what matters."
    ))


def s11_video(prs):
    s = new_slide(prs)
    header(s, 'Walkthrough Demo', 'DEMO VIDEO  |  2-3 minutes')
    txt(s, '(Video submitted as separate file)',
        Inches(0.5), Inches(0.82), Inches(6.0), Inches(0.35),
        sz=13, color=C_DGRAY, italic=True)

    rect(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(4.5),
         RGBColor(15, 22, 35), line=C_TEAL)
    txt(s, '▶', Inches(0.4), Inches(2.6), Inches(12.5), Inches(1.5),
        sz=64, color=C_TEAL, align=PP_ALIGN.CENTER)

    txt(s, 'Key moments in the demo:', Inches(0.4), Inches(6.0), Inches(5.0), Inches(0.38),
        sz=13, bold=True, color=C_LGRAY)
    moments = [
        '1.  Onboarding panel on first load — what ChatFolio is and isn\'t',
        '2.  Conversational profile building — one question at a time',
        '3.  Live sidebar updating in real time as fields are confirmed',
        '4.  Clicking "Generate My Portfolio" — the deliberate pause',
        '5.  Reading the rationale and growth projections',
        '6.  Follow-up adjustment → spinner → portfolio regenerates',
    ]
    bullets(s, moments, Inches(0.4), Inches(6.42), Inches(12.5), Inches(1.0),
            sz=12, prefix='', color=C_LGRAY)

    notes(s, (
        "Play the demo video now — it runs approximately 2 to 3 minutes. "
        "Watch for the live sidebar updates, the generate button moment, and the "
        "follow-up adjustment at the end. After the video, we'll cover lessons learned."
    ))


def s12_lessons(prs):
    s = new_slide(prs)
    header(s, 'Lessons Learned', 'LESSONS LEARNED')
    lessons = [
        (C_TEAL,    'Mixed-initiative interaction works — but only if users feel in control',
         'Users engaged most when they could steer the conversation freely. The edit loop '
         'was used more than we expected — people wanted to iterate, not just accept.'),
        (C_GREEN,   'Financial literacy is a design problem, not a user problem',
         'Every time a user got confused by a term, that was our failure, not theirs. '
         'Inline definitions and plain-language rationale closed the gap more than any other change.'),
        (C_ORANGE,  'Conversation beats forms — but only if it doesn\'t feel like a form in disguise',
         'The one-question-at-a-time approach worked, but users who tried to answer everything '
         'upfront revealed a tension: our sequential design assumed the user would pace with us.'),
        (C_RED,     'The edit loop matters as much as first generation — iteration is the real product',
         'We almost shipped without robust post-generation UX. The think-aloud study made clear '
         'that the moment after "Generate" is where users actually engage with the output.'),
    ]
    for i, (color, title, body) in enumerate(lessons):
        y = Inches(1.1) + i * Inches(1.5)
        rect(s, Inches(0.4), y, Inches(12.5), Inches(1.35), C_CARD)
        rect(s, Inches(0.4), y, Inches(0.12), Inches(1.35), color)
        txt(s, str(i + 1), Inches(0.65), y + Inches(0.1), Inches(0.5), Inches(0.55),
            sz=22, bold=True, color=color)
        txt(s, title, Inches(1.3), y + Inches(0.08), Inches(11.4), Inches(0.45),
            sz=15, bold=True, color=color)
        txt(s, body,  Inches(1.3), y + Inches(0.55), Inches(11.4), Inches(0.72),
            sz=12, color=C_LGRAY)

    notes(s, (
        "Four things we'd tell our past selves. "
        "The financial literacy gap was the most actionable insight — it's fixable with copy, "
        "not architecture. The mixed-initiative finding is broader: users want to be the driver, "
        "not the passenger. And the lesson about the edit loop applies to any generative AI tool — "
        "the output is a starting point for dialogue, not a finished deliverable."
    ))


def s13_thankyou(prs):
    s = new_slide(prs)
    rect(s, 0, 0, W, Inches(0.08), C_TEAL)
    txt(s, 'ChatFolio', Inches(1), Inches(1.6), Inches(11.3), Inches(1.3),
        sz=56, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'Thank You', Inches(1), Inches(2.95), Inches(11.3), Inches(0.7),
        sz=32, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(4.5), Inches(3.78), Inches(4.3), Inches(0.03), C_DGRAY)
    txt(s, 'Chris  ·  Indel  ·  Nate  ·  Matt',
        Inches(1), Inches(3.88), Inches(11.3), Inches(0.5),
        sz=16, color=C_DGRAY, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.8), Inches(4.7), Inches(5.7), Inches(0.6), C_DARK)
    txt(s, 'streamlit run app.py',
        Inches(3.8), Inches(4.72), Inches(5.7), Inches(0.55),
        sz=14, color=C_TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'Questions?', Inches(1), Inches(5.55), Inches(11.3), Inches(0.5),
        sz=20, color=C_LGRAY, align=PP_ALIGN.CENTER)
    notes(s, (
        "Open for questions. The walkthrough video was submitted as a separate file. "
        "The full source code is in the repo — run 'streamlit run app.py' to try it live."
    ))


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    print('ChatFolio Final Presentation Generator')
    print('Building charts...')
    buf_cmp = comparison_chart()

    print('Building slides...')
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s01_title(prs)
    s02_problem(prs)
    s03_solution(prs)
    s04_how_it_works(prs)
    s05_v01_prototype(prs)
    s06_v02_prototype(prs)
    s07_v05_prototype(prs)
    s08_study(prs)
    s09_findings(prs)
    s10_results(prs, buf_cmp)
    s11_video(prs)
    s12_lessons(prs)
    s13_thankyou(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       'chatfolio_final_presentation.pptx')
    prs.save(out)
    print(f'Saved: {out}')
    print(f'Slides: 13  |  Open in PowerPoint to review and export to PDF.')


if __name__ == '__main__':
    main()
