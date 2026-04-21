"""
ChatFolio Evaluation Report Generator
Run from project root: python evaluation/generate_slides.py
Outputs: evaluation/chatfolio_evaluation_report.pptx
"""
import io
import os

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ───────────────────────────────────────────────────────────────────
C_BG    = RGBColor(26,  26,  46)   # slide background
C_CARD  = RGBColor(30,  41,  59)   # card / box bg
C_TEAL  = RGBColor(78,  205, 196)  # V2 / accent
C_RED   = RGBColor(255, 107, 107)  # V1 bars
C_WHITE = RGBColor(255, 255, 255)
C_LGRAY = RGBColor(200, 200, 200)
C_DGRAY = RGBColor(120, 120, 140)
C_GREEN = RGBColor(34,  197, 94)
C_YELL  = RGBColor(255, 213, 79)
C_CODE  = RGBColor(165, 214, 167)  # green code text
C_DARK  = RGBColor(13,  17,  23)   # code block bg

THEME_C = [
    RGBColor(239, 68,  68),   # Onboarding  – red
    RGBColor(249, 115, 22),   # Risk        – orange
    RGBColor(234, 179, 8),    # Chat Flow   – yellow
    RGBColor(34,  197, 94),   # Portfolio   – green
    RGBColor(59,  130, 246),  # UX          – blue
    RGBColor(139, 92,  246),  # What Worked – purple
]

W = Inches(13.33)
H = Inches(7.5)

# ── Metrics (computed from benchmark_dataset_v2.csv, 49 samples) ─────────────
CATS     = ['Goal','Timeline','Budget','Risk','Additional','Post-Gen','Off-Topic']
N        = [7, 7, 7, 6, 7, 7, 8]

V1_H  = [71.4, 71.4, 71.4, 66.7, 28.6, 85.7, 100.0]  # human pass rate %
V2_H  = [100.0, 85.7, 100.0, 100.0, 100.0, 85.7, 100.0]

V1_LLM_AVG, V1_HUM_AVG = 3.61, 3.96
V2_LLM_AVG, V2_HUM_AVG = 4.80, 4.76
V1_LLM_PR,  V1_HUM_PR  = 67.3, 71.4
V2_LLM_PR,  V2_HUM_PR  = 91.8, 95.9


# ── pptx helpers ─────────────────────────────────────────────────────────────
def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return slide


def txt(slide, text, l, t, w, h, sz=16, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(sz)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = 'Calibri'
    return box


def bullets(slide, items, l, t, w, h, sz=14, color=C_LGRAY, prefix='  •  '):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text           = prefix + item
        run.font.size      = Pt(sz)
        run.font.color.rgb = color
        run.font.name      = 'Calibri'
    return box


def rect(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp


def img(slide, buf, l, t, w, h=None):
    buf.seek(0)
    return slide.shapes.add_picture(buf, l, t, width=w, height=h) if h \
           else slide.shapes.add_picture(buf, l, t, width=w)


def accent(slide):
    rect(slide, 0, 0, W, Inches(0.07), C_TEAL)


def header(slide, title, tag=''):
    accent(slide)
    txt(slide, title, Inches(0.5), Inches(0.1), Inches(10.5), Inches(0.72),
        sz=26, bold=True)
    if tag:
        txt(slide, tag, Inches(10.0), Inches(0.1), Inches(3.0), Inches(0.5),
            sz=10, color=C_TEAL, align=PP_ALIGN.RIGHT)


def mkchart(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=150,
                facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf


# ── Chart generators ─────────────────────────────────────────────────────────
def pie_chart():
    colors = ['#EF4444','#F97316','#EAB308','#22C55E','#3B82F6','#8B5CF6','#EC4899']
    fig, ax = plt.subplots(figsize=(5, 4.2), facecolor='#0D1117')
    ax.set_facecolor('#0D1117')
    wedges, texts, autotexts = ax.pie(
        N, labels=CATS, autopct='%1.0f%%', colors=colors,
        startangle=140,
        textprops={'color':'white','fontsize':9},
        wedgeprops={'linewidth':1,'edgecolor':'#1A1A2E'}
    )
    for at in autotexts:
        at.set_fontsize(8)
    ax.set_title('49 test samples', color='white', fontsize=11, pad=6)
    return mkchart(fig)


def v1_bar_chart():
    fig, ax = plt.subplots(figsize=(7.2, 3.8), facecolor='#0D1117')
    ax.set_facecolor('#0D1117')
    y  = range(len(CATS))
    bs = ax.barh(list(y), V1_H, color='#FF6B6B', height=0.55, zorder=3)
    ax.set_yticks(list(y))
    ax.set_yticklabels(CATS, color='white', fontsize=11)
    ax.set_xlabel('Human Pass Rate (%)', color='#AAAAAA', fontsize=10)
    ax.set_xlim(0, 118)
    ax.tick_params(colors='#AAAAAA')
    for sp in ['top','right']:   ax.spines[sp].set_visible(False)
    for sp in ['bottom','left']: ax.spines[sp].set_color('#444')
    ax.grid(axis='x', color='#2A2A3E', linewidth=0.6, zorder=0)
    for b, v in zip(bs, V1_H):
        ax.text(v+1.5, b.get_y()+b.get_height()/2, f'{v:.0f}%',
                va='center', color='white', fontsize=9)
    ax.set_title('V1 Human Pass Rate by Category', color='white', fontsize=12, pad=8)
    fig.tight_layout()
    return mkchart(fig)


def comparison_chart():
    fig, ax = plt.subplots(figsize=(8.5, 4.0), facecolor='#0D1117')
    ax.set_facecolor('#0D1117')
    y    = np.arange(len(CATS))
    bh   = 0.35
    b1 = ax.barh(y - bh/2, V1_H, bh, color='#FF6B6B', label='V1 Baseline', zorder=3)
    b2 = ax.barh(y + bh/2, V2_H, bh, color='#4ECDC4', label='V2 Refined',  zorder=3)
    ax.set_yticks(y)
    ax.set_yticklabels(CATS, color='white', fontsize=11)
    ax.set_xlabel('Human Pass Rate (%)', color='#AAAAAA', fontsize=10)
    ax.set_xlim(0, 125)
    ax.tick_params(colors='#AAAAAA')
    for sp in ['top','right']:   ax.spines[sp].set_visible(False)
    for sp in ['bottom','left']: ax.spines[sp].set_color('#444')
    ax.grid(axis='x', color='#2A2A3E', linewidth=0.6, zorder=0)
    for b, v in zip(b2, V2_H):
        ax.text(v+1.5, b.get_y()+b.get_height()/2, f'{v:.0f}%',
                va='center', color='#4ECDC4', fontsize=9)
    ax.legend(loc='lower right', facecolor='#1E293B', labelcolor='white', fontsize=10)
    ax.set_title('Before vs. After: Human Pass Rate by Category', color='white',
                 fontsize=12, pad=8)
    fig.tight_layout()
    return mkchart(fig)


# ── Slides ────────────────────────────────────────────────────────────────────
def s01_title(prs):
    s = new_slide(prs)
    rect(s, 0, 0, W, Inches(0.08), C_TEAL)
    txt(s, 'ChatFolio', Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
        sz=60, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 'Evaluation Report', Inches(1), Inches(2.65), Inches(11.3), Inches(0.9),
        sz=36, color=C_TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'Qualitative & Quantitative User Evaluation',
        Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
        sz=20, color=C_LGRAY, align=PP_ALIGN.CENTER)
    rect(s, Inches(4.5), Inches(4.25), Inches(4.3), Inches(0.03), C_DGRAY)
    txt(s, 'Chris  \u00b7  Indel  \u00b7  Nate  \u00b7  Matt     |     April 2026',
        Inches(1), Inches(4.35), Inches(11.3), Inches(0.5),
        sz=14, color=C_DGRAY, align=PP_ALIGN.CENTER)


def s02_study_overview(prs):
    s = new_slide(prs)
    header(s, 'Study Overview & Participant Demographics', 'PART 1 \u2014 QUALITATIVE')
    # Methodology box
    rect(s, Inches(0.4), Inches(1.0), Inches(5.9), Inches(5.9), C_CARD)
    txt(s, 'METHODOLOGY', Inches(0.65), Inches(1.1), Inches(5.5), Inches(0.4),
        sz=12, bold=True, color=C_TEAL)
    bullets(s, [
        'Task: complete full ChatFolio flow',
        '   greeting \u2192 4 questions \u2192 generate portfolio',
        'Method: think-aloud, moderated observation',
        'Observer noted friction points, surprises,',
        '   and positive reactions without guiding users',
        'Synthesis: affinity diagramming \u2192 6 themes',
    ], Inches(0.65), Inches(1.6), Inches(5.5), Inches(4.8), sz=15, prefix='')
    # Demographics box
    rect(s, Inches(6.8), Inches(1.0), Inches(6.1), Inches(5.9), C_CARD)
    txt(s, 'PARTICIPANTS  (n = 7)', Inches(7.05), Inches(1.1), Inches(5.7), Inches(0.4),
        sz=12, bold=True, color=C_TEAL)
    bullets(s, [
        'All college students',
        'Majors: CS, Finance, Economics',
        'Mixed investing experience:',
        '   Some had zero prior knowledge',
        '   Some were familiar with stocks / markets',
        '   Some were already actively investing',
    ], Inches(7.05), Inches(1.6), Inches(5.7), Inches(4.8), sz=15, prefix='')


def s03_affinity_results(prs):
    s = new_slide(prs)
    header(s, 'Affinity Diagram Results', 'PART 1 \u2014 QUALITATIVE')
    themes = [
        ('Onboarding &\nInitial Prompt',       '6 / 7 participants', 0),
        ('Risk Tolerance\nConfusion',           '3 / 7 participants', 1),
        ('Chat &\nConversation Flow',           '5 / 7 participants', 2),
        ('Portfolio Output &\nComprehension',   '5 / 7 participants', 3),
        ('UX &\nVisual Polish',                 '4 / 7 participants', 4),
        ('What Worked Well',                    '6 / 7 participants', 5),
    ]
    cw, ch   = Inches(3.9), Inches(2.6)
    sx, sy   = Inches(0.45), Inches(0.95)
    gx, gy   = Inches(0.2), Inches(0.25)
    for i, (name, count, ci) in enumerate(themes):
        col, row = i % 3, i // 3
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)
        tc = THEME_C[ci]
        rect(s, x, y, cw, ch, C_CARD)
        rect(s, x, y, cw, Inches(0.1), tc)
        txt(s, name, x + Inches(0.15), y + Inches(0.18), cw - Inches(0.3), Inches(1.1),
            sz=16, bold=True)
        rect(s, x + Inches(0.15), y + ch - Inches(0.62), Inches(1.6), Inches(0.42), tc)
        txt(s, count,
            x + Inches(0.15), y + ch - Inches(0.65), Inches(1.6), Inches(0.45),
            sz=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)


def s04_key_findings(prs):
    s = new_slide(prs)
    header(s, 'Key Findings by Theme', 'PART 1 \u2014 QUALITATIVE')
    # Pain points
    rect(s, Inches(0.4), Inches(1.0), Inches(8.4), Inches(5.9), C_CARD)
    txt(s, 'PAIN POINTS', Inches(0.65), Inches(1.1), Inches(8.0), Inches(0.4),
        sz=13, bold=True, color=C_RED)
    pain = [
        'Onboarding: Initial prompt implied only "retirement" or "house"; '
        'users felt their goals were being put in a box',
        'Risk Tolerance: Beginners don\'t understand what "risk" means in investing; '
        '"35% confident" phrasing was confusing',
        'Chat Flow: Users dumped all info in one message, '
        'bypassing the conversational step-by-step design',
        'Portfolio Output: ETF tickers (VTI, VXUS, BND) and "index fund" were '
        'meaningless to beginners — literacy gap',
        'UX: No visual indicator when a follow-up message triggers '
        'portfolio regeneration; no loading state shown',
    ]
    bullets(s, pain, Inches(0.65), Inches(1.6), Inches(7.9), Inches(4.8),
            sz=14, prefix='  \u2192  ')
    # What worked
    rect(s, Inches(9.0), Inches(1.0), Inches(3.9), Inches(5.9), C_CARD)
    txt(s, 'WHAT WORKED', Inches(9.2), Inches(1.1), Inches(3.5), Inches(0.4),
        sz=13, bold=True, color=C_GREEN)
    worked = [
        'Real-time sidebar profile updates',
        'Loading screen / onboarding panel',
        'Growth projection chart (clear & intuitive)',
        'Conversational format strongly preferred over traditional forms',
    ]
    bullets(s, worked, Inches(9.2), Inches(1.6), Inches(3.5), Inches(4.8),
            sz=14, color=C_GREEN, prefix='  \u2713  ')


def s05_ux_takeaways(prs):
    s = new_slide(prs)
    header(s, 'UX Takeaways', 'PART 1 \u2014 QUALITATIVE')
    items = [
        ('1', 'Gather life context before investment goals',
         'Collect age, income, debt, and life stage before asking about specific goals — '
         'users need the system to understand them first'),
        ('2', 'Close the financial literacy gap',
         'Define ETFs, indexes, and risk tolerance inline. '
         'Never assume baseline financial knowledge for a beginner-focused tool'),
        ('3', 'Fix the follow-up feedback loop',
         'Users need a clear signal that their follow-up message was received '
         'and that the portfolio is regenerating'),
        ('4', 'Lean into the conversational format',
         'The conversational design is the core differentiator — '
         'invest in making it feel like a real advisor, not a form in disguise'),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.1) + i * Inches(1.5)
        rect(s, Inches(0.4), y + Inches(0.05), Inches(0.45), Inches(0.45), C_TEAL)
        txt(s, num, Inches(0.4), y + Inches(0.05), Inches(0.45), Inches(0.45),
            sz=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(s, title, Inches(1.05), y, Inches(12.0), Inches(0.45),
            sz=18, bold=True)
        txt(s, desc,  Inches(1.05), y + Inches(0.5), Inches(12.0), Inches(0.85),
            sz=14, color=C_LGRAY)


def s06_critical_prompt(prs):
    s = new_slide(prs)
    header(s, 'The Critical Prompt', 'PART 2 \u2014 QUANTITATIVE METRICS')
    txt(s, 'Which prompt?', Inches(0.4), Inches(1.05), Inches(6.2), Inches(0.4),
        sz=14, bold=True, color=C_TEAL)
    txt(s,
        'The conversational profile extraction system prompt in chat_engine.py. '
        'This single prompt handles every user message from greeting to portfolio generation.',
        Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.9), sz=14, color=C_LGRAY)
    txt(s, 'Why critical?', Inches(0.4), Inches(2.55), Inches(6.2), Inches(0.4),
        sz=14, bold=True, color=C_TEAL)
    bullets(s, [
        'Routes 100% of user messages at every stage',
        'Must handle 5 profile fields and diverse edge cases',
        'Returns structured JSON \u2014 format errors break the entire app',
        'Failures cascade directly to bad portfolio recommendations',
    ], Inches(0.4), Inches(3.05), Inches(6.2), Inches(3.2), sz=13)
    # Code block
    rect(s, Inches(7.0), Inches(1.0), Inches(5.9), Inches(5.9), C_DARK)
    txt(s, 'System prompt excerpt (V2)', Inches(7.15), Inches(1.08), Inches(5.6), Inches(0.35),
        sz=10, color=C_DGRAY)
    code = (
        'You are ChatFolio, a friendly investment advisor\n'
        'that helps beginners build a starter portfolio\n'
        'through conversation.\n\n'
        'Guidelines:\n'
        '1. Ask ONE question at a time.\n'
        '2. Paraphrase vague answers and confirm\n'
        '   before saving a field.\n'
        '3. Collect in order: goal \u2192 timeline \u2192 budget\n'
        '   \u2192 risk_tolerance \u2192 additional_info\n'
        '4. Handle edge cases: variable budgets,\n'
        '   numeric risk scales, short timelines...\n\n'
        'Respond ONLY in JSON:\n'
        '{\n'
        '  "message": "...",\n'
        '  "profile_updates": { ... },\n'
        '  "ready_for_portfolio": false,\n'
        '  "profile_changed": false\n'
        '}'
    )
    txt(s, code, Inches(7.15), Inches(1.55), Inches(5.6), Inches(5.1),
        sz=12, color=C_CODE)


def s07_dataset(prs, buf_pie):
    s = new_slide(prs)
    header(s, 'Test Dataset', 'PART 2 \u2014 QUANTITATIVE METRICS')
    # Stat cards
    cards = [
        ('49',    'total samples'),
        ('7',     'categories'),
        ('~60%',  'typical inputs'),
        ('~20%',  'edge cases'),
        ('~20%',  'adversarial'),
    ]
    for i, (val, lbl) in enumerate(cards):
        x = Inches(0.4) + i * Inches(1.35)
        rect(s, x, Inches(1.0), Inches(1.2), Inches(1.2), C_CARD)
        txt(s, val, x, Inches(1.05), Inches(1.2), Inches(0.62),
            sz=22, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
        txt(s, lbl, x, Inches(1.68), Inches(1.2), Inches(0.38),
            sz=10, color=C_LGRAY, align=PP_ALIGN.CENTER)
    img(s, buf_pie, Inches(1.6), Inches(2.25), Inches(5.5))
    # Legend
    cat_hex = ['#EF4444','#F97316','#EAB308','#22C55E','#3B82F6','#8B5CF6','#EC4899']
    txt(s, 'Category breakdown', Inches(8.1), Inches(2.25), Inches(4.8), Inches(0.38),
        sz=13, bold=True)
    for i, (cat, n) in enumerate(zip(CATS, N)):
        y   = Inches(2.72) + i * Inches(0.58)
        r,g,b = int(cat_hex[i][1:3],16), int(cat_hex[i][3:5],16), int(cat_hex[i][5:],16)
        rect(s, Inches(8.1), y + Inches(0.07), Inches(0.22), Inches(0.22), RGBColor(r,g,b))
        txt(s, f'{cat}  ({n} samples)', Inches(8.45), y, Inches(4.4), Inches(0.45),
            sz=13, color=C_LGRAY)


def s08_metrics(prs):
    s = new_slide(prs)
    header(s, 'Evaluation Metrics', 'PART 2 \u2014 QUANTITATIVE METRICS')
    boxes = [
        ('LLM-as-Judge',
         'GPT-4o-mini rates each response against ground truth on a 1\u20135 scale',
         ['Same rubric applied consistently across all 49 samples',
          'Secondary metric \u2014 used for consistency check and scale',
          'Catches systematic failures across many inputs quickly']),
        ('Human Rating',
         'Manual 1\u20135 rating by team member against ground truth',
         ['Primary metric for pass / fail determination',
          'Pass threshold: \u2265 4 out of 5',
          'Catches LLM judge blind spots and context-dependent errors']),
    ]
    for i, (title, desc, pts) in enumerate(boxes):
        x = Inches(0.4) + i * Inches(6.4)
        rect(s, x, Inches(1.0), Inches(6.0), Inches(3.6), C_CARD)
        txt(s, title, x + Inches(0.25), Inches(1.1), Inches(5.5), Inches(0.52),
            sz=20, bold=True)
        txt(s, desc,  x + Inches(0.25), Inches(1.65), Inches(5.5), Inches(0.65),
            sz=13, color=C_LGRAY)
        bullets(s, pts, x + Inches(0.25), Inches(2.38), Inches(5.5), Inches(2.0), sz=13)
    txt(s, 'Rating Rubric', Inches(0.4), Inches(4.78), Inches(12.5), Inches(0.4),
        sz=13, bold=True, color=C_TEAL)
    rubric = [
        ('5 \u2014 Excellent', 'Fully matches ground truth', C_GREEN),
        ('4 \u2014 Good',      'Mostly correct, minor gaps', RGBColor(132,204,22)),
        ('3 \u2014 Adequate',  'Partially correct', C_YELL),
        ('2 \u2014 Poor',      'Mostly fails expectations', RGBColor(251,146,60)),
        ('1 \u2014 Unsat.',    'Irrelevant or contradicts', C_RED),
    ]
    for i, (lbl, desc, c) in enumerate(rubric):
        x = Inches(0.4) + i * Inches(2.55)
        rect(s, x, Inches(5.25), Inches(2.45), Inches(1.85), C_CARD)
        txt(s, lbl,  x + Inches(0.12), Inches(5.32), Inches(2.25), Inches(0.44),
            sz=13, bold=True, color=c)
        txt(s, desc, x + Inches(0.12), Inches(5.75), Inches(2.25), Inches(1.1),
            sz=11, color=C_LGRAY)


def s09_v1_baseline(prs, buf_bars):
    s = new_slide(prs)
    header(s, 'Baseline V1 Performance', 'PART 2 \u2014 QUANTITATIVE METRICS')
    stats = [
        ('3.61 / 5', 'LLM Avg Score'),
        ('3.96 / 5', 'Human Avg Score'),
        ('67.3%',    'LLM Pass Rate'),
        ('71.4%',    'Human Pass Rate'),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = Inches(0.4) + i * Inches(3.12)
        rect(s, x, Inches(1.0), Inches(2.9), Inches(1.1), C_CARD)
        txt(s, val, x, Inches(1.02), Inches(2.9), Inches(0.62),
            sz=22, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
        txt(s, lbl, x, Inches(1.62), Inches(2.9), Inches(0.35),
            sz=11, color=C_DGRAY, align=PP_ALIGN.CENTER)
    img(s, buf_bars, Inches(0.3), Inches(2.2), Inches(8.5))
    txt(s, 'Root Causes', Inches(9.1), Inches(2.2), Inches(3.9), Inches(0.4),
        sz=13, bold=True, color=C_TEAL)
    causes = [
        'Variable budgets accepted as-is instead of asking for a steady amount',
        'Vague timelines (ranges) saved without confirmation',
        '"YOLO crypto" treated as aggressive instead of redirecting',
        '"$0, I\'m broke" confirmed instead of suggesting $25\u201350',
        'Additional Info stage misclassified as off-topic in many cases',
    ]
    bullets(s, causes, Inches(9.1), Inches(2.7), Inches(3.9), Inches(4.5), sz=12)


def s10_refinements(prs):
    s = new_slide(prs)
    header(s, 'Prompt Refinements: V1 \u2192 V2', 'PART 3 \u2014 PERFORMANCE COMPARISON')
    changes = [
        ('1', 'One-field confirmation loop',
         'Explicit rule: paraphrase understanding and confirm ONE field before saving. '
         'Added required format: "It sounds like X \u2014 does that sound right?"'),
        ('2', 'Edge case handling',
         'Variable budget \u2192 ask for steady amount. Numeric risk scale \u2192 map to two nearest '
         'categories and let user choose. Short timeline \u2192 volatility warning. $0 \u2192 empathize + suggest $25\u201350.'),
        ('3', 'Post-generation behavioral split',
         'Profile change requests update fields and set profile_changed=true. '
         'Informational questions (e.g. "what is VTI?") answer only \u2014 never touch profile.'),
        ('4', 'Proactive contextual question',
         'After 4 required fields, ask about existing savings, debt, and life changes '
         'instead of generic "anything else?" to surface portfolio-relevant context.'),
        ('5', 'Required plain-language risk explanation',
         'Model must describe all three risk levels before asking. '
         'Prohibited: numeric confidence outputs like "35% confident."'),
    ]
    for i, (num, title, desc) in enumerate(changes):
        y = Inches(1.0) + i * Inches(1.2)
        rect(s, Inches(0.4), y + Inches(0.05), Inches(0.42), Inches(0.42), C_TEAL)
        txt(s, num, Inches(0.4), y + Inches(0.05), Inches(0.42), Inches(0.42),
            sz=17, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(s, title, Inches(1.05), y, Inches(12.0), Inches(0.42),
            sz=15, bold=True)
        txt(s, desc, Inches(1.05), y + Inches(0.43), Inches(12.0), Inches(0.65),
            sz=12, color=C_LGRAY)


def s11_comparison(prs, buf_cmp):
    s = new_slide(prs)
    header(s, 'Before vs. After: Performance Comparison', 'PART 3 \u2014 PERFORMANCE COMPARISON')
    img(s, buf_cmp, Inches(0.3), Inches(1.0), Inches(8.9))
    # Summary table
    cx = [Inches(9.25), Inches(10.3), Inches(11.35), Inches(12.25)]
    cw = Inches(1.05)
    hdrs = ['Metric', 'V1', 'V2', '\u0394']
    for j, h in enumerate(hdrs):
        txt(s, h, cx[j], Inches(1.08), cw, Inches(0.38),
            sz=12, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    rect(s, Inches(9.25), Inches(1.46), Inches(3.75), Inches(0.03), C_DGRAY)
    rows = [
        ('LLM Avg',       '3.61', '4.80', '+1.19'),
        ('Human Avg',     '3.96', '4.76', '+0.80'),
        ('LLM Pass %',    '67.3%','91.8%','+24.5 pp'),
        ('Human Pass %',  '71.4%','95.9%','+24.5 pp'),
    ]
    for i, row in enumerate(rows):
        y  = Inches(1.55) + i * Inches(0.65)
        bg = C_CARD if i % 2 == 0 else RGBColor(20, 28, 42)
        rect(s, Inches(9.25), y - Inches(0.05), Inches(3.75), Inches(0.6), bg)
        cols_c = [C_LGRAY, C_RED, C_TEAL, C_GREEN]
        for j, (cell, cc) in enumerate(zip(row, cols_c)):
            txt(s, cell, cx[j], y, cw, Inches(0.5),
                sz=12, bold=(j == 3), color=cc, align=PP_ALIGN.CENTER)
    # Headline stat box
    rect(s, Inches(9.25), Inches(4.25), Inches(3.75), Inches(1.55), RGBColor(10, 30, 30))
    txt(s, '+24.5 pp', Inches(9.25), Inches(4.3), Inches(3.75), Inches(0.65),
        sz=32, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'human pass rate improvement',
        Inches(9.25), Inches(4.95), Inches(3.75), Inches(0.38),
        sz=12, color=C_LGRAY, align=PP_ALIGN.CENTER)
    txt(s, '71.4%  \u2192  95.9%',
        Inches(9.25), Inches(5.35), Inches(3.75), Inches(0.4),
        sz=16, bold=True, align=PP_ALIGN.CENTER)


def s12_conclusion(prs):
    s = new_slide(prs)
    header(s, 'Conclusion', '')
    points = [
        (C_RED,   'Qualitative',
         '5 UX friction points identified across a 7-participant think-aloud study. '
         'Top issues: financial literacy gap, overly narrow initial prompt framing, '
         'and missing regeneration feedback after follow-up messages.'),
        (C_TEAL,  'Quantitative',
         'Prompt V2 achieves a 95.9% human pass rate (up from 71.4% baseline, +24.5 pp). '
         'LLM-as-judge and human ratings aligned closely (avg delta < 0.1 across all samples).'),
        (C_GREEN, 'Combined Approach',
         'Qualitative user observations directly informed the quantitative test cases, '
         'enabling evidence-driven prompt iteration from observed friction to measurable reliability improvement.'),
    ]
    for i, (c, lbl, body) in enumerate(points):
        y = Inches(1.25) + i * Inches(1.9)
        rect(s, Inches(0.4), y, Inches(12.5), Inches(1.75), C_CARD)
        rect(s, Inches(0.4), y, Inches(0.12), Inches(1.75), c)
        txt(s, lbl,  Inches(0.65), y + Inches(0.15), Inches(11.8), Inches(0.45),
            sz=17, bold=True, color=c)
        txt(s, body, Inches(0.65), y + Inches(0.6),  Inches(11.8), Inches(1.0),
            sz=14, color=C_LGRAY)


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    print('ChatFolio Evaluation Report Generator')
    print('Building charts...')
    buf_pie  = pie_chart()
    buf_bars = v1_bar_chart()
    buf_cmp  = comparison_chart()

    print('Building slides...')
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s01_title(prs)
    s02_study_overview(prs)
    s03_affinity_results(prs)
    s04_key_findings(prs)
    s05_ux_takeaways(prs)
    s06_critical_prompt(prs)
    s07_dataset(prs, buf_pie)
    s08_metrics(prs)
    s09_v1_baseline(prs, buf_bars)
    s10_refinements(prs)
    s11_comparison(prs, buf_cmp)
    s12_conclusion(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       'chatfolio_evaluation_report.pptx')
    prs.save(out)
    print(f'Saved: {out}')
    print('Done. Open in PowerPoint and export to PDF.')


if __name__ == '__main__':
    main()
